from pptx import Presentation
from pptx.util import Inches, Pt
import sys

def create_ppt():
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "VoxLog"
    subtitle.text = "Intelligent Voice Recognition Attendance & Academic Management System\n\nPresenter: [Your Name/Team Name]\nProject Context: Next-Generation Classroom Automation"

    slides_data = [
        {
            "title": "The Problem - Traditional Attendance",
            "bullets": [
                "Time-Consuming: Manual roll calls waste valuable instructional time (10-15 minutes per class).",
                "Proxy Attendance: Traditional methods are highly prone to buddy-punching or fake proxies.",
                "Administrative Burden: Manual entry of attendance data into digital systems is tedious and error-prone."
            ]
        },
        {
            "title": "The Problem - Academic Management",
            "bullets": [
                "Scattered Systems: Faculty leaves, class substitutions, and scheduling are managed on paper/multitude of software.",
                "Lack of Real-Time Data: Administrators lack immediate access to live attendance statistics and academic progress.",
                "Inefficient Workflows: Handling sudden faculty absences and rearranging timetables manually causes campus chaos."
            ]
        },
        {
            "title": "The Solution - Introduction to VoxLog",
            "bullets": [
                "What is VoxLog?: An intelligent system blending voice recognition with automated verification for fast, accurate attendance.",
                "Unified Platform: A single centralized web application handling attendance, academic scheduling, and faculty management.",
                "Roles & Portals: Dedicated, secure dashboards for both System Administrators and Faculty members."
            ]
        },
        {
            "title": "Solution Feature - Voice-Based Attendance",
            "bullets": [
                "Fast & Accurate Verification: Uses voice biometrics to instantly verify student presence.",
                "Fraud Prevention: Eliminates proxy attendance since each student's voice print is unique and impossible to share.",
                "Seamless Process: Allows teachers to take attendance rapidly without disrupting the flow of the class."
            ]
        },
        {
            "title": "Solution Feature - Dedicated Faculty Portal",
            "bullets": [
                "Class & Progress Tracking: Easy access to 'My Class', gradebooks, and student academic progress.",
                "Automated Timetables: Real-time schedule viewing that updates dynamically.",
                "Leave & Substitution Workflow: Faculty can initiate leave requests and coordinate substitute teachers effortlessly."
            ]
        },
        {
            "title": "Solution Feature - Admin Dashboard",
            "bullets": [
                "Centralized Management: Easily oversee departments, courses, programs, students, and user access.",
                "Real-time Monitoring: Access to a Master Attendance Calendar and live tracking of all school activities.",
                "Streamlined Approvals: One-click management for faculty leave approvals and automated substitutions."
            ]
        },
        {
            "title": "Technology Stack",
            "bullets": [
                "Frontend: React.js powered by Vite for a lighting-fast, responsive interface.",
                "Backend Database: Firebase (Firestore) for scalable, real-time NoSQL data syncing.",
                "Authentication: Secure user login & role-management via Firebase Auth.",
                "Biometrics Integration: Voice recognition algorithms to process and verify audio inputs."
            ]
        },
        {
            "title": "Key Benefits & Impact",
            "bullets": [
                "Maximized Teaching Time: Reclaims hours of instructional time over the semester.",
                "High Security: Biometric voice data ensures 100% reliable attendance logs.",
                "Operational Efficiency: Drastically reduces paperwork for administrative tasks.",
                "Eco-Friendly: Transitions institutions into paperless, digital-first environments."
            ]
        },
        {
            "title": "Conclusion & Future Scope",
            "bullets": [
                "Conclusion: VoxLog modernizes classroom management, proving attendance can be secure, quick, and integrated.",
                "Future Scope: Integration with Parent SMS notifications.",
                "Future Scope: Advanced AI analytics to predict student performance based on attendance.",
                "Future Scope: Multi-language voice support for diverse classrooms."
            ]
        }
    ]

    bullet_slide_layout = prs.slide_layouts[1]

    for sd in slides_data:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = sd["title"]

        tf = body_shape.text_frame
        for i, bullet in enumerate(sd["bullets"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    prs.save("VoxLog_Presentation.pptx")
    print("Presentation created successfully as VoxLog_Presentation.pptx")

if __name__ == '__main__':
    create_ppt()
